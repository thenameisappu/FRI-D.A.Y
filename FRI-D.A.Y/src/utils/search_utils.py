import os
import string
import shutil
import sqlite3
import datetime 
from docx import Document
import PyPDF2
import openpyxl
import pptx



 # --- Document Content Search Logic ---
def get_folders_to_scan(scan_type="targeted"):
    if scan_type == "full":
        drives = [f"{L}:\\" for L in string.ascii_uppercase if os.path.exists(f"{L}:\\")]
        return drives
    else:
        home = os.path.expanduser("~")
        folders = [os.path.join(home, d) for d in ("Documents", "Desktop", "Downloads")]
        return [f for f in folders if os.path.exists(f)]
def read_txt(fp): return open(fp, "r", encoding="utf-8", errors="ignore").read()
def read_pdf(fp): return "\n".join(p.extract_text() for p in PyPDF2.PdfReader(fp).pages if p.extract_text())
def read_docx(fp): return "\n".join(p.text for p in Document(fp).paragraphs)
def read_xlsx(fp): return "\n".join(str(c.value) for ws in openpyxl.load_workbook(fp, read_only=True) for r in ws.iter_rows() for c in r if c.value)
def read_pptx(fp): return "\n".join(s.text for slide in pptx.Presentation(fp).slides for s in slide.shapes if hasattr(s, "text"))
FILE_READERS = {".txt": read_txt, ".pdf": read_pdf, ".docx": read_docx, ".xlsx": read_xlsx, ".pptx": read_pptx}
def search_documents(keyword, folders):
    matches = []
    for folder in folders:
        print(f"[*] Scanning: {folder}")
        for root, _, files in os.walk(folder):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in FILE_READERS:
                    filepath = os.path.join(root, file)
                    try:
                        if keyword.lower() in FILE_READERS[ext](filepath).lower():
                            matches.append({"file_name": file, "file_path": filepath})
                            print(f"  [+] Match: {filepath}")
                    except Exception:
                        continue
    return matches
def handle_document_search(speak, listen_command, get_text_command):
    speak("What keyword should I search for in your documents?")
    keyword = listen_command() or get_text_command()
    if not keyword:
        speak("No keyword provided. Cancelling search.")
        return
    speak("Should I do a fast, targeted scan of common folders, or a full system scan which is very slow?")
    response = listen_command() or get_text_command()
    scan_type = "full" if response and "full" in response else "targeted"
  
    speak(f"Okay, starting a {scan_type} scan for the keyword '{keyword}'. This might take some time.")
    folders = get_folders_to_scan(scan_type)
    results = search_documents(keyword, folders)
    if not results:
        speak(f"I couldn't find any documents containing the keyword '{keyword}'.")
    else:
        speak(f"I found {len(results)} matching documents.")
        for i, match in enumerate(results[:5]):
            speak(f"Match {i+1} is {match['file_name']}")
        if len(results) > 5:
            speak(f"There are {len(results) - 5} more results. Check the console for the full list.")


# --- Browser History Search Logic ---

BROWSER_PATHS = {
    "Chrome": r"~\AppData\Local\Google\Chrome\User Data\Default\History",
    "Edge": r"~\AppData\Local\Microsoft\Edge\User Data\Default\History",
    "Firefox": r"~\AppData\Roaming\Mozilla\Firefox\Profiles",
    "Brave": r"~\AppData\Local\BraveSoftware\Brave-Browser\User Data\Default\History",
    "Opera": r"~\AppData\Roaming\Opera Software\Opera Stable\History"
}

def search_chromium_history(path, keyword):
    results = []
    temp_db = "temp_history.db"
    try:
        shutil.copy2(path, temp_db)
        conn = sqlite3.connect(temp_db)
        cursor = conn.cursor()
        query = "SELECT title, last_visit_time FROM urls WHERE url LIKE ? OR title LIKE ? ORDER BY last_visit_time DESC"
        cursor.execute(query, (f"%{keyword}%", f"%{keyword}%"))
        for title, ts in cursor.fetchall():
            dt = datetime.datetime(1601, 1, 1) + datetime.timedelta(microseconds=ts)
            results.append({"title": title, "last_visited": dt.strftime("%Y-%m-%d %H:%M")})
        conn.close()
        os.remove(temp_db)
    except Exception:
        pass
    return results

def search_firefox_history(path, keyword):
    # Firefox history search
    results = []
    for profile in os.listdir(path):
        places_db = os.path.join(path, profile, "places.sqlite")
        if os.path.exists(places_db):
            # ... (logic similar to chromium search)
            pass
    return results

def search_all_history(keyword):
    """Searches all installed browsers for a keyword."""
    combined = []
    installed = {name: os.path.expanduser(p) for name, p in BROWSER_PATHS.items() if os.path.exists(os.path.expanduser(p))}
    
    for browser, path in installed.items():
        if browser == "Firefox":
            results = search_firefox_history(path, keyword)
        else:
            results = search_chromium_history(path, keyword)
        for r in results:
            r["browser"] = browser
        combined.extend(results)
    return combined

def handle_history_search(speak, listen_command, get_text_command):
    """Voice-driven workflow for searching browser history."""
    speak("What keyword should I search for in your browser history?")
    keyword = listen_command() or get_text_command()
    if not keyword:
        speak("No keyword provided. Cancelling search.")
        return
        
    speak(f"Okay, searching all browser histories for '{keyword}'.")
    results = search_all_history(keyword)
    
    if not results:
        speak(f"I couldn't find anything in your browser history for '{keyword}'.")
    else:
        speak(f"I found {len(results)} items in your history.")
        for i, res in enumerate(results[:5]): # Read out top 5
            speak(f"From {res['browser']}, you visited a page titled {res['title']} on {res['last_visited']}")
        if len(results) > 5:
            speak("There are more results. Check the console for the full list.")